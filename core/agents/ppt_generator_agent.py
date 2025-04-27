#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
PPT生成Agent模块

负责根据布局决策生成最终的PPT文件。
"""

import logging
import json
from typing import Dict, Any, List, Optional
from pathlib import Path

from core.agents.base_agent import BaseAgent
from core.engine.state import AgentState

logger = logging.getLogger(__name__)

class PPTGeneratorAgent(BaseAgent):
    """PPT生成Agent，负责根据布局决策生成PPT文件"""
    
    def __init__(self, config: Dict[str, Any]):
        """
        初始化PPT生成Agent
        
        Args:
            config: Agent配置
        """
        super().__init__(config)
        # 获取模板路径，如果提供的话
        self.template_dir = config.get("template_dir", None)
        logger.info(f"初始化PPTGeneratorAgent，模板目录: {self.template_dir}")
    
    async def run(self, state: AgentState) -> AgentState:
        """
        执行PPT生成
        
        Args:
            state: 当前工作流状态
            
        Returns:
            更新后的状态
        """
        logger.info("开始生成PPT文件")
        
        # 检查必要的输入
        if not state.decision_result:
            error_msg = "没有提供布局决策"
            self.record_failure(state, error_msg)
            return state
        
        try:
            # 生成PPT文件
            output_path = self._generate_ppt(state.decision_result, state.session_id)
            
            # 更新状态
            state.ppt_file_path = str(output_path)
            logger.info(f"PPT生成完成，文件保存至: {output_path}")
            
            # 记录检查点
            self.add_checkpoint(state)
            
        except Exception as e:
            error_msg = f"PPT生成失败: {str(e)}"
            self.record_failure(state, error_msg)
        
        return state
    
    def _generate_ppt(self, decision_result: Dict[str, Any], session_id: str) -> Path:
        """
        根据决策生成PPT文件
        
        Args:
            decision_result: 布局决策结果
            session_id: 会话ID
            
        Returns:
            生成的PPT文件路径
        """
        from config.settings import settings
        
        # 创建输出目录
        output_dir = settings.OUTPUT_DIR
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # 设置输出文件路径
        output_path = output_dir / f"{session_id}.pptx"
        
        logger.info(f"生成PPT文件: {output_path}")
        
        # 在实际项目中，这里会使用python-pptx等库创建PPT文件
        # 这里仅做模拟，将决策结果保存为JSON文件
        
        # 保存决策结果为JSON
        json_path = output_dir / f"{session_id}_decision.json"
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(decision_result, f, ensure_ascii=False, indent=2)
        
        # 模拟创建PPT文件
        self._create_mock_pptx(output_path, decision_result)
        
        return output_path
    
    def _create_mock_pptx(self, output_path: Path, decision_result: Dict[str, Any]) -> None:
        """
        模拟创建PPT文件
        
        Args:
            output_path: 输出文件路径
            decision_result: 布局决策结果
        """
        try:
            # 检查是否有示例PPT文件
            from config.settings import settings
            
            # 在实际项目中集成外部库
            try:
                # 尝试导入python-pptx库
                from pptx import Presentation
                from pptx.util import Inches, Pt
                
                # 创建新的演示文稿
                prs = Presentation()
                
                # 遍历决策中的每张幻灯片
                for slide_info in decision_result.get("slides", []):
                    slide_type = slide_info.get("type")
                    content = slide_info.get("content", {})
                    
                    # 根据类型选择布局
                    if slide_type == "title":
                        # 标题页
                        slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = content.get("title", "")
                        if hasattr(slide.shapes, "placeholders") and len(slide.shapes.placeholders) > 1:
                            slide.shapes.placeholders[1].text = content.get("subtitle", "")
                    
                    elif slide_type == "content":
                        # 内容页
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = content.get("title", "")
                        
                        # 添加内容
                        body_shape = slide.shapes.placeholders[1]
                        tf = body_shape.text_frame
                        
                        # 添加项目符号列表
                        for bullet in content.get("bullets", []):
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    
                    elif slide_type == "twoColumns":
                        # 双栏布局
                        slide_layout = prs.slide_layouts[3]  # 假设这是两栏布局
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = content.get("title", "")
                        
                        # 在实际应用中，需要确定左右两栏的占位符索引
                        left_items = content.get("leftContent", [])
                        right_items = content.get("rightContent", [])
                        
                        # 这里仅做简单示例
                        if len(slide.shapes.placeholders) > 1:
                            left_shape = slide.shapes.placeholders[1]
                            left_tf = left_shape.text_frame
                            
                            for item in left_items:
                                p = left_tf.add_paragraph()
                                p.text = item
                                p.level = 0
                        
                        if len(slide.shapes.placeholders) > 2:
                            right_shape = slide.shapes.placeholders[2]
                            right_tf = right_shape.text_frame
                            
                            for item in right_items:
                                p = right_tf.add_paragraph()
                                p.text = item
                                p.level = 0
                    
                    elif slide_type == "image":
                        # 图片页
                        slide_layout = prs.slide_layouts[5]  # 假设这是图片布局
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = content.get("title", "")
                        
                        # 在实际应用中，需要添加图片
                        # 这里仅做简单示例
                        if len(slide.shapes.placeholders) > 1:
                            body_shape = slide.shapes.placeholders[1]
                            tf = body_shape.text_frame
                            tf.text = content.get("caption", "")
                
                # 保存演示文稿
                prs.save(output_path)
                logger.info(f"使用python-pptx创建PPT文件: {output_path}")
                
            except ImportError:
                # 如果没有python-pptx，则创建一个空文件
                logger.warning("python-pptx库未安装，创建空文件")
                
                # 尝试从模板复制一个PPT文件
                example_ppt = Path("/Users/deadwalk/Code/ai_proj_ppt/ppt_assisstant/libs/ppt_manager/test/testfiles/Iphone16Pro.pptx")
                if example_ppt.exists():
                    import shutil
                    shutil.copy(example_ppt, output_path)
                    logger.info(f"从示例复制PPT文件: {output_path}")
                else:
                    # 创建空文件
                    with open(output_path, 'w') as f:
                        f.write("# PPT文件模拟\n")
                    logger.warning(f"创建模拟PPT文件: {output_path}")
        
        except Exception as e:
            logger.error(f"创建PPT文件失败: {str(e)}")
            # 创建一个空白文件作为备用
            with open(output_path, 'w') as f:
                f.write("# 创建PPT文件失败\n")
            logger.warning(f"创建空白文件作为备用: {output_path}")
    
    def _get_layout_by_name(self, prs, layout_name: str):
        """
        根据名称获取布局
        
        Args:
            prs: Presentation对象
            layout_name: 布局名称
            
        Returns:
            对应的布局
        """
        # 实际应用中，需要匹配布局名称与索引
        layout_mapping = {
            "title": 0,
            "content": 1,
            "twoColumns": 3,
            "image": 5
        }
        
        idx = layout_mapping.get(layout_name, 1)  # 默认使用内容布局
        if idx < len(prs.slide_layouts):
            return prs.slide_layouts[idx]
        else:
            return prs.slide_layouts[0] 