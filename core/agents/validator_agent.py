#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
验证Agent模块

负责验证生成的PPT是否符合要求，包括完整性、格式和内容覆盖度。
"""

import logging
import os
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path

from core.agents.base_agent import BaseAgent
from core.engine.state import AgentState

logger = logging.getLogger(__name__)

class ValidatorAgent(BaseAgent):
    """验证Agent，负责验证生成的PPT文件"""
    
    def __init__(self, config: Dict[str, Any]):
        """
        初始化验证Agent
        
        Args:
            config: Agent配置
        """
        super().__init__(config)
        # 配置验证规则
        self.validation_rules = config.get("validation_rules", {})
        logger.info(f"初始化ValidatorAgent，验证规则数量: {len(self.validation_rules)}")
    
    async def run(self, state: AgentState) -> AgentState:
        """
        执行PPT验证
        
        Args:
            state: 当前工作流状态
            
        Returns:
            更新后的状态
        """
        logger.info("开始验证PPT文件")
        
        # 增加验证尝试次数
        state.validation_attempts += 1
        
        # 检查必要的输入
        if not state.ppt_file_path:
            error_msg = "没有提供PPT文件路径"
            self.record_failure(state, error_msg)
            return state
        
        # 检查文件是否存在
        ppt_path = Path(state.ppt_file_path)
        if not ppt_path.exists():
            error_msg = f"PPT文件不存在: {ppt_path}"
            self.record_failure(state, error_msg)
            return state
        
        try:
            # 验证PPT文件
            validation_result, issues = self._validate_ppt(
                ppt_path, 
                state.content_structure, 
                state.decision_result
            )
            
            # 如果有验证问题，记录它们
            if not validation_result and issues:
                for issue in issues:
                    self.record_failure(state, f"验证问题: {issue}")
                
                logger.warning(f"PPT验证失败，发现{len(issues)}个问题")
            else:
                logger.info("PPT验证通过")
            
            # 记录检查点
            self.add_checkpoint(state)
            
        except Exception as e:
            error_msg = f"PPT验证过程出错: {str(e)}"
            self.record_failure(state, error_msg)
        
        return state
    
    def _validate_ppt(self, ppt_path: Path, 
                    content_structure: Dict[str, Any],
                    decision_result: Dict[str, Any]) -> Tuple[bool, List[str]]:
        """
        验证PPT文件
        
        Args:
            ppt_path: PPT文件路径
            content_structure: 内容结构
            decision_result: 布局决策结果
            
        Returns:
            验证结果元组 (是否通过, 问题列表)
        """
        logger.info(f"验证PPT文件: {ppt_path}")
        
        # 问题列表
        issues = []
        
        # 检查文件是否存在
        if not ppt_path.exists():
            issues.append(f"文件不存在: {ppt_path}")
            return False, issues
        
        # 检查文件大小
        file_size = os.path.getsize(ppt_path)
        if file_size < 100:  # 如果文件太小，可能是空文件
            issues.append(f"文件大小异常 ({file_size} 字节)，可能是空文件")
        
        # 在实际项目中，这里会使用python-pptx来解析PPT并验证内容
        try:
            # 尝试导入python-pptx
            from pptx import Presentation
            
            # 打开PPT文件
            prs = Presentation(ppt_path)
            
            # 检查幻灯片数量
            expected_slides = len(decision_result.get("slides", []))
            actual_slides = len(prs.slides)
            
            if actual_slides < expected_slides:
                issues.append(f"幻灯片数量不足，预期: {expected_slides}，实际: {actual_slides}")
            
            # 检查标题是否存在
            if content_structure and content_structure.get("title"):
                title_found = False
                for slide in prs.slides:
                    if hasattr(slide.shapes, "title") and slide.shapes.title:
                        if content_structure["title"] in slide.shapes.title.text:
                            title_found = True
                            break
                
                if not title_found:
                    issues.append(f"未找到标题: {content_structure['title']}")
            
            # 验证其他内容（在实际项目中会有更复杂的验证）
            
        except ImportError:
            logger.warning("python-pptx库未安装，使用备用验证方法")
            
            # 备用验证：检查JSON决策文件是否存在
            json_path = ppt_path.with_suffix("").with_suffix(".json")
            if not json_path.exists():
                issues.append(f"缺少决策JSON文件: {json_path}")
        
        except Exception as e:
            issues.append(f"验证时发生错误: {str(e)}")
        
        # 返回验证结果
        passed = len(issues) == 0
        return passed, issues 